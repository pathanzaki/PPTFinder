from flask import Flask, request, send_file, render_template, jsonify, send_from_directory
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json, io, os, re, uuid
import requests

app = Flask(__name__)

# ── PUT YOUR GROQ API KEY HERE ──────────────────────────
API_KEY = os.environ.get("GROQ_API_KEY")
# ────────────────────────────────────────────────────────

SITES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_sites")
os.makedirs(SITES_DIR, exist_ok=True)

SW = Inches(13.33)
SH = Inches(7.5)

# ── helpers ──────────────────────────────────────────────
def C(r, g, b):
    return RGBColor(r, g, b)

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def tb(slide, x, y, w, h, text, sz,
       bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    if color is None:
        color = C(0xFF, 0xFF, 0xFF)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = str(text)
    r.font.size      = Pt(sz)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    r.font.name      = "Calibri"
    return box

def set_bg(slide, color):
    b = slide.background.fill
    b.solid()
    b.fore_color.rgb = color

# ── Topic-aware color themes ──────────────────────────────
def get_theme(topic):
    t = topic.lower()
    if any(w in t for w in ['food','restaurant','pizza','cafe','chef','cuisine','cook']):
        return dict(
            dark=C(0x18,0x06,0x02), mid=C(0x2A,0x0E,0x05),
            light=C(0xFF,0xF5,0xEE), panel=C(0x20,0x09,0x03),
            a1=C(0xE8,0x5D,0x04), a2=C(0xFF,0xBA,0x08), a3=C(0xD0,0x00,0x00)
        )
    if any(w in t for w in ['nature','climate','green','eco','forest','environment','sustainab']):
        return dict(
            dark=C(0x04,0x12,0x07), mid=C(0x08,0x21,0x0D),
            light=C(0xF0,0xFA,0xF1), panel=C(0x06,0x18,0x09),
            a1=C(0x2D,0x9E,0x4F), a2=C(0x74,0xC6,0x9D), a3=C(0x40,0x91,0x6C)
        )
    if any(w in t for w in ['health','medical','hospital','doctor','fitness','wellness','body','mental']):
        return dict(
            dark=C(0x03,0x0C,0x1E), mid=C(0x05,0x16,0x30),
            light=C(0xF0,0xF6,0xFF), panel=C(0x04,0x10,0x26),
            a1=C(0x00,0x77,0xCC), a2=C(0x00,0xC2,0xFF), a3=C(0x38,0x5E,0xF5)
        )
    if any(w in t for w in ['finance','money','bank','invest','stock','economy','wealth','trading']):
        return dict(
            dark=C(0x04,0x10,0x06), mid=C(0x08,0x1C,0x0A),
            light=C(0xF2,0xFA,0xF3), panel=C(0x06,0x14,0x08),
            a1=C(0x00,0x8F,0x39), a2=C(0x34,0xD3,0x89), a3=C(0x06,0x6E,0x29)
        )
    if any(w in t for w in ['art','design','fashion','photo','music','creative','film','cinema']):
        return dict(
            dark=C(0x10,0x04,0x18), mid=C(0x1C,0x07,0x2A),
            light=C(0xFB,0xF0,0xFF), panel=C(0x14,0x05,0x1E),
            a1=C(0xAA,0x00,0xFF), a2=C(0xFF,0x5C,0xE1), a3=C(0x7B,0x2F,0xBE)
        )
    if any(w in t for w in ['travel','tourism','hotel','adventure','explore','destination']):
        return dict(
            dark=C(0x02,0x0C,0x18), mid=C(0x03,0x16,0x2A),
            light=C(0xF0,0xF8,0xFF), panel=C(0x03,0x10,0x20),
            a1=C(0xFF,0x6B,0x35), a2=C(0xFF,0xC3,0x47), a3=C(0x06,0x8F,0xBF)
        )
    if any(w in t for w in ['space','astro','nasa','planet','cosmos','universe','rocket']):
        return dict(
            dark=C(0x02,0x02,0x12), mid=C(0x05,0x05,0x22),
            light=C(0xF2,0xF2,0xFF), panel=C(0x04,0x04,0x1C),
            a1=C(0x7B,0x2F,0xFF), a2=C(0x00,0xE5,0xFF), a3=C(0xFF,0x66,0x00)
        )
    # default: deep tech / professional
    return dict(
        dark=C(0x06,0x06,0x14), mid=C(0x0E,0x0E,0x24),
        light=C(0xF3,0xF4,0xFF), panel=C(0x10,0x10,0x2A),
        a1=C(0x7C,0x3A,0xFF), a2=C(0x00,0xD4,0xFF), a3=C(0xFF,0x6B,0x35)
    )

# ── SLIDE BUILDERS ────────────────────────────────────────

def s_title(prs, d, T):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, T['dark'])

    # Background decorative blocks
    rect(sl, 6.8,  -0.8, 7.2, 5.2, T['a1'])
    rect(sl, 7.8,  -0.8, 6.2, 5.2, T['dark'])
    rect(sl, 8.8,  -0.3, 5.5, 4.5, T['mid'])

    # Bottom stripe (two tone)
    rect(sl, 0,    7.05, 8.5,  0.45, T['a1'])
    rect(sl, 8.5,  7.05, 4.83, 0.45, T['a2'])

    # Left vertical bar
    rect(sl, 0,    0,    0.30, 7.5,  T['a2'])

    # Thin accent line
    rect(sl, 0.30, 3.1,  9.0,  0.06, T['a2'])

    # Title
    tb(sl, 0.65, 1.2, 9.8, 1.9,
       d['title'], 52, bold=True, color=C(0xFF, 0xFF, 0xFF))

    # Explanation (first 2 sentences only for title slide)
    expl = d.get('explanation', '')
    if expl:
        sentences = [s.strip() for s in expl.split('.') if s.strip()]
        subtitle = '. '.join(sentences[:2]) + '.' if sentences else expl
        tb(sl, 0.65, 3.25, 9.2, 1.3,
           subtitle, 17, italic=True, color=C(0xCC, 0xCC, 0xEE))

    # Bullet pills row
    bullets = d.get('bullets', [])
    bx = 0.65
    for b in bullets[:4]:
        pill_w = min(len(b) * 0.115 + 0.5, 3.8)
        rect(sl, bx, 5.7, pill_w, 0.48, T['panel'])
        tb(sl, bx + 0.14, 5.7, pill_w - 0.2, 0.48, b, 11.5, color=T['a2'])
        bx += pill_w + 0.18
        if bx > 12.5:
            break

    # Footer
    tb(sl, 0.65, 7.1, 9, 0.35,
       "PPTFinders AI  |  pptfinders.com",
       10, italic=True, color=C(0xFF, 0xFF, 0xFF))


def s_two_col(prs, d, T, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, T['light'])

    # Left dark panel
    rect(sl, 0,    0,    6.0,  7.5,  T['mid'])
    rect(sl, 0,    0,    6.0,  0.24, T['a1'])
    rect(sl, 0,    0,    0.24, 7.5,  T['a2'])
    rect(sl, 0,    7.26, 6.0,  0.24, T['a2'])

    # Right top thin line
    rect(sl, 6.0,  0,    7.33, 0.24, C(0xE0, 0xE4, 0xFF))

    # Slide number
    tb(sl, 12.55, 0.22, 0.7, 0.4,
       str(num), 11, color=T['a2'], align=PP_ALIGN.RIGHT)

    # Title (left panel)
    tb(sl, 0.38, 0.44, 5.25, 1.55,
       d['title'], 26, bold=True, color=C(0xFF, 0xFF, 0xFF))

    # Divider line
    rect(sl, 0.38, 2.1, 5.0, 0.05, T['a2'])

    # Full explanation paragraph
    expl = d.get('explanation', '')
    if expl:
        tb(sl, 0.38, 2.25, 5.3, 4.8,
           expl, 13, color=C(0xCC, 0xCC, 0xDD))

    # Right: numbered bullet cards
    colors = [T['a1'], T['a2'], T['a3'],
              C(0x28, 0x9B, 0x6C), C(0xEE, 0x55, 0x22)]
    cy = 0.48
    for i, b in enumerate(d.get('bullets', [])[:5]):
        cc = colors[i % len(colors)]
        # card shadow
        rect(sl, 6.28, cy + 0.06, 6.85, 1.15, C(0xD5, 0xD8, 0xF5))
        # card bg
        rect(sl, 6.22, cy,        6.85, 1.15, C(0xFF, 0xFF, 0xFF))
        # left color stripe
        rect(sl, 6.22, cy,        0.15, 1.15, cc)
        # number circle
        circ = sl.shapes.add_shape(
            9, Inches(6.52), Inches(cy + 0.28), Inches(0.6), Inches(0.6))
        circ.fill.solid()
        circ.fill.fore_color.rgb = cc
        circ.line.fill.background()
        tb(sl, 6.52, cy + 0.28, 0.6, 0.6,
           str(i + 1), 15, bold=True,
           color=C(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        # bullet text
        tb(sl, 7.28, cy + 0.1, 5.6, 0.98,
           b, 13, color=C(0x10, 0x10, 0x24))
        cy += 1.28

    tb(sl, 6.22, 7.1, 6, 0.3,
       "PPTFinders AI", 9,
       color=C(0x88, 0x88, 0xAA), align=PP_ALIGN.RIGHT)


def s_fulltext(prs, d, T, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, T['dark'])

    # Top accent bar
    rect(sl, 0,    0,    13.33, 0.18, T['a1'])
    # Right decorative column
    rect(sl, 11.5, 0.18, 1.83,  7.32, T['panel'])
    rect(sl, 12.95, 0,   0.38,  7.5,  T['a2'])

    # Slide num
    tb(sl, 11.65, 0.3, 1.1, 0.55,
       f"/{num:02d}", 14, color=T['a2'], align=PP_ALIGN.CENTER)

    # Title
    tb(sl, 0.55, 0.35, 10.7, 1.1,
       d['title'], 33, bold=True, color=C(0xFF, 0xFF, 0xFF))

    # Full explanation - large and prominent
    expl = d.get('explanation', '')
    if expl:
        tb(sl, 0.55, 1.6, 10.6, 2.75,
           expl, 15.5, color=C(0xCC, 0xCC, 0xEE))

    # KEY POINTS section
    bullets = d.get('bullets', [])
    if bullets:
        rect(sl, 0,    4.55, 13.33, 0.055, T['a2'])
        tb(sl, 0.55, 4.72, 5, 0.38,
           "▸  KEY POINTS", 9, bold=True, color=T['a2'])

        cx = [0.45, 3.9, 7.35]
        bcolors = [T['a1'], T['a2'], T['a3']]
        for i, b in enumerate(bullets[:3]):
            rect(sl, cx[i], 5.22, 3.1, 2.0, T['panel'])
            rect(sl, cx[i], 5.22, 3.1, 0.14, bcolors[i])
            circ = sl.shapes.add_shape(
                9, Inches(cx[i] + 0.2), Inches(5.46),
                Inches(0.52), Inches(0.52))
            circ.fill.solid()
            circ.fill.fore_color.rgb = bcolors[i]
            circ.line.fill.background()
            tb(sl, cx[i] + 0.2, 5.46, 0.52, 0.52,
               "★", 11, color=C(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
            tb(sl, cx[i] + 0.86, 5.38, 2.08, 1.75,
               b, 12.5, color=C(0xDD, 0xDD, 0xFF))

    tb(sl, 0.55, 7.12, 8, 0.3,
       "PPTFinders AI", 9,
       italic=True, color=C(0x55, 0x55, 0x88))


def s_callout(prs, d, T, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, T['light'])

    # Header band
    rect(sl, 0,    0,    13.33, 1.45, T['mid'])
    rect(sl, 0,    0,    0.28,  7.5,  T['a1'])
    rect(sl, 0.28, 0,    0.14,  1.45, T['a2'])
    rect(sl, 0.28, 1.45, 13.05, 0.055, T['a2'])

    tb(sl, 0.65, 0.24, 11.5, 0.98,
       d['title'], 30, bold=True, color=C(0xFF, 0xFF, 0xFF))
    tb(sl, 12.55, 0.28, 0.7, 0.4,
       str(num), 11, color=T['a2'], align=PP_ALIGN.RIGHT)

    # Explanation left-bottom
    expl = d.get('explanation', '')
    if expl:
        tb(sl, 0.55, 1.65, 8.5, 2.65,
           expl, 14.5, color=C(0x11, 0x11, 0x33))

    # Remaining bullets strip
    all_b = d.get('bullets', [])
    if len(all_b) > 4:
        rect(sl, 0.28, 4.45, 8.6, 0.055, T['a1'])
        bx = 0.55
        for b in all_b[4:7]:
            tb(sl, bx, 4.62, 2.7, 0.85,
               f"→  {b}", 12.5, color=C(0x22, 0x22, 0x44))
            bx += 2.85

    # Right: 4 coloured callout boxes
    colors = [T['a1'], T['a2'], T['a3'], C(0x28, 0x9B, 0x6C)]
    by = 1.52
    for i, b in enumerate(all_b[:4]):
        cc = colors[i % 4]
        # shadow
        rect(sl, 9.18, by + 0.07, 3.92, 1.38, C(0xCC, 0xCE, 0xEE))
        rect(sl, 9.1,  by,        3.92, 1.38, cc)
        tb(sl, 9.22, by + 0.06, 0.75, 0.56,
           f"{i+1:02d}", 24, bold=True, color=C(0xFF, 0xFF, 0xFF))
        rect(sl, 9.22, by + 0.68, 3.62, 0.05, C(0xFF, 0xFF, 0xFF))
        tb(sl, 9.22, by + 0.76, 3.62, 0.6,
           b, 12.5, color=C(0xFF, 0xFF, 0xFF))
        by += 1.5

    tb(sl, 0.55, 7.12, 6, 0.3,
       "PPTFinders AI", 9,
       italic=True, color=C(0x88, 0x88, 0xAA))


def s_timeline(prs, d, T, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, T['dark'])

    rect(sl, 0, 0, 13.33, 0.18, T['a1'])

    tb(sl, 0.55, 0.32, 11.8, 1.0,
       d['title'], 32, bold=True, color=C(0xFF, 0xFF, 0xFF))
    tb(sl, 12.55, 0.28, 0.7, 0.4,
       str(num), 11, color=T['a2'], align=PP_ALIGN.RIGHT)

    expl = d.get('explanation', '')
    if expl:
        tb(sl, 0.55, 1.45, 12.2, 1.22,
           expl, 14.5, color=C(0xBB, 0xBB, 0xDD))

    # Horizontal timeline bar
    rect(sl, 0.55, 3.42, 12.2, 0.12, T['a2'])

    bullets = d.get('bullets', [])
    n  = max(len(bullets), 1)
    sw = 12.2 / n
    colors = [T['a1'], T['a2'], T['a3'],
              C(0x28, 0x9B, 0x6C), C(0xC0, 0x30, 0xFF)]

    for i, b in enumerate(bullets[:5]):
        cx   = 0.55 + i * sw + sw / 2 - 0.36
        cc   = colors[i % len(colors)]

        # connector line above bar
        rect(sl, cx + 0.36, 2.8, 0.07, 0.65, T['a2'])

        # circle node
        circ = sl.shapes.add_shape(
            9, Inches(cx), Inches(2.45), Inches(0.74), Inches(0.74))
        circ.fill.solid()
        circ.fill.fore_color.rgb = cc
        circ.line.fill.background()
        tb(sl, cx, 2.45, 0.74, 0.74,
           str(i + 1), 16, bold=True,
           color=C(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

        # card below bar
        card_x = max(0.3, cx - 0.48)
        card_w = min(sw - 0.14, 2.28)
        rect(sl, card_x, 3.7, card_w, 3.0, T['panel'])
        rect(sl, card_x, 3.7, card_w, 0.15, cc)
        tb(sl, card_x + 0.14, 3.95, card_w - 0.24, 2.68,
           b, 11.5, color=C(0xDD, 0xDD, 0xFF))

    tb(sl, 0.55, 7.12, 8, 0.3,
       "PPTFinders AI", 9,
       italic=True, color=C(0x55, 0x55, 0x88))


def s_grid(prs, d, T, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, T['light'])

    rect(sl, 0,    0,    13.33, 1.35, T['mid'])
    rect(sl, 0,    0,    0.28,  7.5,  T['a2'])
    rect(sl, 0.28, 1.35, 13.05, 0.055, T['a1'])

    tb(sl, 0.55, 0.24, 12.0, 0.88,
       d['title'], 30, bold=True, color=C(0xFF, 0xFF, 0xFF))
    tb(sl, 12.55, 0.26, 0.7, 0.4,
       str(num), 11, color=T['a2'], align=PP_ALIGN.RIGHT)

    bullets = d.get('bullets', [])
    pos = [
        (0.42, 1.55), (4.55, 1.55), (8.68, 1.55),
        (0.42, 4.10), (4.55, 4.10), (8.68, 4.10),
    ]
    colors = [T['a1'], T['a2'], T['a3'],
              C(0x28, 0x9B, 0x6C), C(0xEE, 0x55, 0x22), C(0x44, 0x44, 0xCC)]

    for i, b in enumerate(bullets[:6]):
        if i >= len(pos):
            break
        px, py = pos[i]
        cc = colors[i % len(colors)]
        # shadow
        rect(sl, px + 0.07, py + 0.07, 3.82, 2.3, C(0xCC, 0xCE, 0xEE))
        # card
        rect(sl, px, py, 3.82, 2.3, C(0xFF, 0xFF, 0xFF))
        # top color band
        rect(sl, px, py, 3.82, 0.27, cc)
        # icon circle
        circ = sl.shapes.add_shape(
            9, Inches(px + 0.14), Inches(py + 0.44),
            Inches(0.60), Inches(0.60))
        circ.fill.solid()
        circ.fill.fore_color.rgb = cc
        circ.line.fill.background()
        tb(sl, px + 0.14, py + 0.44, 0.60, 0.60,
           str(i + 1), 15, bold=True,
           color=C(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        # text
        tb(sl, px + 0.90, py + 0.36, 2.76, 1.82,
           b, 12.5, color=C(0x10, 0x10, 0x22))

    tb(sl, 0.42, 7.12, 8, 0.3,
       "PPTFinders AI", 9,
       italic=True, color=C(0x88, 0x88, 0xAA))


def s_conclusion(prs, d, T):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, T['dark'])

    rect(sl, 0, 0,    13.33, 0.20, T['a2'])
    rect(sl, 0, 7.30, 13.33, 0.20, T['a1'])
    rect(sl, 0, 0,    0.30,  7.5,  T['a1'])
    rect(sl, 13.03, 0, 0.30, 7.5,  T['a2'])

    # Centered content box
    rect(sl, 1.2, 1.2, 10.93, 5.0, T['panel'])
    rect(sl, 1.2, 1.2, 10.93, 0.20, T['a1'])
    rect(sl, 1.2, 6.0, 10.93, 0.20, T['a2'])

    # Big title
    tb(sl, 1.6, 1.5, 10.13, 1.65,
       d['title'], 46, bold=True,
       color=C(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # Divider
    rect(sl, 5.16, 3.3, 2.99, 0.07, T['a2'])

    # Summary text
    expl = d.get('explanation', '')
    if expl:
        tb(sl, 1.6, 3.5, 10.13, 1.6,
           expl, 15.5, color=C(0xBB, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

    # Bottom key takeaway pills
    bullets = d.get('bullets', [])
    if bullets:
        bx    = 1.5
        bstep = 10.33 / max(len(bullets[:4]), 1)
        for b in bullets[:4]:
            tb(sl, bx, 5.2, bstep - 0.1, 0.62,
               f"◆  {b}", 12.5, color=T['a2'], align=PP_ALIGN.CENTER)
            bx += bstep

    tb(sl, 1.2, 7.05, 10.93, 0.35,
       "Generated by PPTFinders AI  |  pptfinders.com",
       10, italic=True, color=C(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)


# ── Layout cycle ──────────────────────────────────────────
LAYOUTS = [
    "two_col",  "fulltext", "callout",  "two_col",
    "timeline", "grid",     "fulltext", "callout",
    "two_col",  "grid",     "timeline", "fulltext",
    "callout",  "grid",     "two_col",  "timeline",
    "fulltext", "grid",     "callout",  "two_col",
]

def build_pptx(slides, topic):
    T   = get_theme(topic)
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    num = 1
    for i, sd in enumerate(slides):
        st = sd.get('slide_type', 'content')
        if st == 'title':
            s_title(prs, sd, T)
        elif st == 'conclusion':
            s_conclusion(prs, sd, T)
        else:
            lay = LAYOUTS[(i - 1) % len(LAYOUTS)]
            if   lay == 'two_col':  s_two_col(prs, sd, T, num)
            elif lay == 'fulltext': s_fulltext(prs, sd, T, num)
            elif lay == 'callout':  s_callout(prs, sd, T, num)
            elif lay == 'timeline': s_timeline(prs, sd, T, num)
            elif lay == 'grid':     s_grid(prs, sd, T, num)
            num += 1
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

def clean_json(raw):
    # Remove markdown
    if "```" in raw:
        parts = raw.split("```")
        for p in parts:
            p = p.strip()
            if p.startswith("json"):
                p = p[4:]
            if p.startswith("["):
                raw = p
                break

    # Extract JSON array
    start = raw.find("[")
    end = raw.rfind("]") + 1

    if start != -1 and end > start:
        raw = raw[start:end]

    return raw.strip()
# ── Groq PPT content ─────────────────────────────────────
def gen_ppt_content(prompt, num_slides):
    num_slides = max(5, min(30, int(num_slides)))
    cc = num_slides - 2
    system = f"""You are a world-class presentation writer and subject-matter expert.
Return ONLY a raw JSON array of exactly {num_slides} slide objects. No markdown, no backticks, no commentary.

Each object MUST have ALL 4 fields:
  "title"       : punchy specific title, max 9 words
  "slide_type"  : "title" | "content" | "conclusion"
  "explanation" : EXACTLY 4-5 full sentences of expert-level content:
                  • Sentence 1: clear definition / context
                  • Sentences 2-3: specific facts, statistics, real examples, mechanisms
                  • Sentence 4-5: implications, applications, or future outlook
                  Write like a Forbes article or university textbook. NO vague generalities.
  "bullets"     : list of EXACTLY 5 strings, each 10-18 words, specific and factual.
                  Include numbers/percentages where relevant. Not vague one-liners.

RULES:
  Slide 1 → "title", Slide {num_slides} → "conclusion"
  Slides 2-{num_slides-1} → "content" (exactly {cc} slides, each covering a DIFFERENT subtopic)
  Return ONLY the raw JSON array."""

    client = Groq(api_key=API_KEY)
    resp = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {"role": "system", "content": system},
            {"role": "user",   "content": f"Create a {num_slides}-slide presentation: {prompt}"}
        ],
        temperature=0.58,
        max_tokens=5000,
    )
    raw = resp.choices[0].message.content.strip()
    raw = clean_json(raw)

try:
    slides = json.loads(raw)
except Exception as e:
    print("RAW RESPONSE:\n", raw)
    raise ValueError("Invalid JSON from AI")
    if slides:
        slides[0]["slide_type"]  = "title"
        slides[-1]["slide_type"] = "conclusion"
    return slides


# ── Groq Website generation ───────────────────────────────
WEBSITE_PROMPT = """You are a senior full-stack developer and award-winning UI/UX designer.
Generate a COMPLETE, STUNNING, fully-responsive single-page website.

Return ONLY valid JSON (no markdown):
{"site_title": "...", "description": "...", "html": "...complete HTML string..."}

THE HTML MUST INCLUDE ALL OF THESE:

SECTIONS:
1. Sticky NAVBAR — logo, 5 links, CTA button, working mobile hamburger
2. HERO — full viewport, animated gradient headline, subtext, 2 CTA buttons, decorative CSS shapes
3. ABOUT — 2-column (text + stats grid: 4 big numbers)
4. SERVICES/FEATURES — 6-card grid with emoji icons, titles, descriptions, hover lift
5. STATS — dark-band section with 4 animated count-up counters
6. TESTIMONIALS — 3 glassmorphism cards with quote, name, role, CSS star rating
7. FAQ — 5 Q&A items with smooth accordion (JS max-height animation)
8. CONTACT — form (name, email, message, submit) + success state
9. FOOTER — 3 columns (logo+desc, links, contact), social icons, copyright

CSS REQUIREMENTS (inside <style>):
  - :root variables for all colors (theme must match topic)
  - Google Fonts: 2 fonts (link in <head>)
  - Mobile-first, breakpoints at 768px and 1100px
  - Animated hero gradient background (@keyframes)
  - Hero headline: gradient text (background-clip:text)
  - Card hover: translateY(-8px) + shadow deepens
  - Glassmorphism: backdrop-filter:blur(12px) + semi-transparent bg
  - Custom scrollbar (webkit)
  - .reveal class: opacity:0 translateY(25px) → visible: opacity:1 translateY(0)
  - Consistent spacing scale (8/16/24/32/48/64/96px)
  - Button: gradient bg, border-radius:50px, hover scale(1.04)

JAVASCRIPT (inside <script>, NO external libs):
  - Hamburger toggle with X animation
  - Navbar shrink + shadow after 80px scroll
  - Smooth scroll on all anchor links
  - IntersectionObserver → add .visible to .reveal elements
  - Count-up animation on stats (0 → target, triggered by observer)
  - FAQ accordion using max-height toggle
  - Active nav link highlighting by scroll position
  - Form submit: preventDefault, validate, show success message

CONTENT:
  - All text must be specific to the topic in the prompt
  - Real business name from prompt, industry-specific copy
  - Realistic stats, testimonial names, FAQ questions

COLOR: warm palette for food/lifestyle, cool/blue for tech, green for eco/health, etc.

Return ONLY the JSON object. No markdown outside it."""

def gen_website(prompt):
    client = Groq(api_key=API_KEY)
    resp = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {"role": "system", "content": WEBSITE_PROMPT},
            {"role": "user",   "content": f"Build a complete website for: {prompt}"}
        ],
        temperature=0.68,
        max_tokens=5000,
    )
    raw = resp.choices[0].message.content.strip()
    # Strip fences
    if "```" in raw:
        for part in raw.split("```"):
            part = part.strip()
            if part.startswith("json"):
                part = part[4:].strip()
            if part.startswith("{"):
                raw = part
                break
    s = raw.find("{")
    e = raw.rfind("}") + 1
    if s != -1 and e > s:
        raw = raw[s:e]
    return json.loads(raw)


# ── Flask routes ──────────────────────────────────────────
@app.route("/")
def index():
    return render_template("index.html")


@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    data   = request.get_json(force=True)
    prompt = data.get("prompt", "").strip()
    n      = int(data.get("num_slides", 12))

    if not prompt:
        return jsonify({"error": "Please enter a topic."}), 400
    if not (5 <= n <= 30):
        return jsonify({"error": "Slide count must be 5–30."}), 400

    try:
        slides    = gen_ppt_content(prompt, n)
        pptx_data = build_pptx(slides, prompt)
        fname     = re.sub(r"[^a-z0-9_]", "_", prompt[:35].lower()) + ".pptx"
        return send_file(
            io.BytesIO(pptx_data),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=fname,
        )
    except json.JSONDecodeError:
        return jsonify({"error": "AI returned bad JSON. Try again."}), 500
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


@app.route("/generate-website", methods=["POST"])
def generate_website():
    data   = request.get_json(force=True)
    prompt = data.get("prompt", "").strip()

    if not prompt:
        return jsonify({"error": "Please describe your website."}), 400

    try:
        result    = gen_website(prompt)
        html_code = result.get("html", "")
        if not html_code:
            return jsonify({"error": "No HTML returned. Try again."}), 500

        site_id  = str(uuid.uuid4())[:8]
        safe     = re.sub(r"[^a-z0-9_]", "_", prompt[:28].lower())
        filename = f"{safe}_{site_id}.html"
        fpath    = os.path.join(SITES_DIR, filename)
        with open(fpath, "w", encoding="utf-8") as f:
            f.write(html_code)

        return jsonify({
            "success":     True,
            "html":        html_code,
            "site_title":  result.get("site_title", "Your Website"),
            "description": result.get("description", ""),
            "filename":    filename,
        })
    except json.JSONDecodeError:
        return jsonify({"error": "AI returned bad JSON. Try again."}), 500
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


@app.route("/preview/<filename>")
def preview(filename):
    return send_from_directory(SITES_DIR, filename)


@app.route("/download-site/<filename>")
def download_site(filename):
    return send_from_directory(SITES_DIR, filename, as_attachment=True)


if __name__ == "__main__":
    print("\n✦ PPTFinders AI Studio → http://127.0.0.1:5000\n")
    app.run(debug=True)
